import logging
from collections import Counter
from datetime import datetime, timezone
import io
import re
from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from scraper import fetch_channel_info, fetch_videos, enrich_videos_with_engagement, parse_upload_date, parse_duration, format_duration, parse_number

app = Flask(__name__)

# configure simple logging for the app
logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(name)s: %(message)s')
logger = logging.getLogger(__name__)

PRIMARY_COLOR = RGBColor(12, 48, 97)
ACCENT_COLOR = RGBColor(0, 128, 179)
BACKGROUND_COLOR = RGBColor(245, 248, 251)
HEADER_TEXT_COLOR = RGBColor(255, 255, 255)
BODY_TEXT_COLOR = RGBColor(45, 62, 80)
SECONDARY_TEXT_COLOR = RGBColor(80, 96, 112)

STOPWORDS = {
    'about', 'after', 'again', 'against', 'among', 'around', 'because', 'before',
    'being', 'between', 'both', 'could', 'every', 'first', 'for', 'from', 'have',
    'however', 'into', 'least', 'many', 'more', 'most', 'other', 'over', 'their',
    'there', 'these', 'this', 'those', 'through', 'under', 'which', 'while', 'with',
    'your', 'years', 'video', 'videos', 'company', 'channel', 'official', 'new',
    'youtube', 'media', 'content', 'view', 'views', 'watch', 'best', 'most'
}

def compute_cadence(dates):
    dates = sorted([d for d in dates if d])
    if not dates:
        return {
            'cadence': 'Unknown',
            'average_days': None,
            'videos_per_month': 0,
            'latest_upload': 'Unknown',
        }
    latest = dates[-1].strftime('%Y-%m-%d')
    if len(dates) == 1:
        return {
            'cadence': 'Single video sample',
            'average_days': None,
            'videos_per_month': 0.5,
            'latest_upload': latest,
        }
    intervals = [(dates[i] - dates[i - 1]).days for i in range(1, len(dates))]
    average_days = sum(intervals) / len(intervals)
    if average_days <= 7:
        cadence = 'Weekly'
    elif average_days <= 14:
        cadence = 'Bi-weekly'
    elif average_days <= 30:
        cadence = 'Monthly'
    else:
        cadence = 'Irregular'
    videos_per_month = round(len(dates) * 30 / max((dates[-1] - dates[0]).days, 1), 1)
    return {
        'cadence': cadence,
        'average_days': average_days,
        'videos_per_month': videos_per_month,
        'latest_upload': latest,
    }


def extract_topic_keywords(texts, top_n=8):
    counter = Counter()
    for text in texts:
        if not text or not isinstance(text, str):
            continue
        words = re.findall(r"[a-zA-Z0-9]{4,}", text.lower())
        for word in words:
            if word in STOPWORDS or word.isdigit():
                continue
            counter[word] += 1
    return [word for word, _ in counter.most_common(top_n)]


def extract_theme_phrases(texts, top_n=6):
    counter = Counter()
    for text in texts:
        if not text or not isinstance(text, str):
            continue
        tokens = [word for word in re.findall(r"[a-zA-Z0-9]{4,}", text.lower()) if word not in STOPWORDS]
        for first, second in zip(tokens, tokens[1:]):
            if first != second:
                counter[f"{first} {second}"] += 1
    return [phrase for phrase, _ in counter.most_common(top_n)]



def compute_company_summary(company_name):
    channel = fetch_channel_info(company_name)
    videos = fetch_videos(company_name, limit=8)
    videos = enrich_videos_with_engagement(videos, max_details=6)
    if not channel:
        channel = {
            'title': company_name,
            'channel_id': None,
            'description': 'No verified channel data found.',
            'subscribers': 'Unknown',
            'subscriber_count': 0,
            'view_count': 'Unknown',
            'total_views': sum(v['view_count'] for v in videos),
            'video_count': len(videos),
            'channel_url': '',
        }

    valid_views = [v['view_count'] for v in videos if v['view_count'] is not None]
    valid_likes = [v['like_count'] for v in videos if isinstance(v.get('like_count'), int)]
    valid_comments = [v['comment_count'] for v in videos if isinstance(v.get('comment_count'), int)]
    upload_dates = [parse_upload_date(v['upload_date']) for v in videos if v.get('upload_date')]
    texts = [f"{v['title']} {v['description']} {' '.join(v.get('tags', []))}" for v in videos]
    topic_keywords = extract_topic_keywords(texts)
    theme_phrases = extract_theme_phrases(texts)
    cadence = compute_cadence(upload_dates)

    stats = {
        'total_videos': len(videos),
        'average_views': int(sum(valid_views) / max(len(valid_views), 1)) if valid_views else 0,
        'total_views': sum(valid_views),
        'average_likes': int(sum(valid_likes) / max(len(valid_likes), 1)) if valid_likes else 0,
        'average_comments': int(sum(valid_comments) / max(len(valid_comments), 1)) if valid_comments else 0,
        'average_duration_seconds': int(sum(v['duration_seconds'] for v in videos) / max(len(videos), 1)) if videos else 0,
        'most_viewed_video': max(videos, key=lambda v: v['view_count'] or 0) if videos else None,
        'most_liked_video': max(videos, key=lambda v: v.get('like_count') or 0) if videos else None,
        'most_commented_video': max(videos, key=lambda v: v.get('comment_count') or 0) if videos else None,
        'top_videos': sorted(videos, key=lambda v: v['view_count'] or 0, reverse=True)[:5],
        'top_topics': topic_keywords,
        'top_phrases': theme_phrases,
        'content_gaps': [],
        'video_cadence': cadence,
        'videos': videos,
    }
    return {
        'company_name': company_name,
        'channel': channel,
        'stats': stats,
    }


def score_companies(summaries):
    max_views = max((summary['stats']['average_views'] for summary in summaries), default=1)
    max_subs = max((summary.get('channel', {}).get('subscriber_count') or 0 for summary in summaries), default=1)
    max_top_video = max((summary['stats']['most_viewed_video']['view_count'] if summary['stats']['most_viewed_video'] else 0 for summary in summaries), default=1)

    def cadence_score(cadence_label):
        mapping = {
            'Weekly': 1.0,
            'Bi-weekly': 0.85,
            'Monthly': 0.7,
            'Irregular': 0.45,
            'Single video sample': 0.35,
            'Unknown': 0.4,
        }
        return mapping.get(cadence_label, 0.4)

    for summary in summaries:
        views_norm = summary['stats']['average_views'] / max_views if max_views else 0
        subs_norm = (summary.get('channel', {}).get('subscriber_count') or 0) / max_subs if max_subs else 0
        top_norm = (summary['stats']['most_viewed_video']['view_count'] / max_top_video) if summary['stats']['most_viewed_video'] and max_top_video else 0
        cadence_norm = cadence_score(summary['stats']['video_cadence']['cadence'])
        score = min(1.0, 0.45 * views_norm + 0.3 * subs_norm + 0.15 * top_norm + 0.1 * cadence_norm)
        summary['score'] = int(score * 100)

    ranked = sorted(summaries, key=lambda s: s['score'], reverse=True)
    for rank, summary in enumerate(ranked, start=1):
        summary['rank'] = rank
    return ranked


def build_insights(summaries):
    insights = []
    if summaries:
        sorted_by_score = sorted(summaries, key=lambda s: s['score'], reverse=True)
        leader = sorted_by_score[0]
        insights.append(f"Leading channel: {leader['company_name']} with a score of {leader['score']}.")
        insights.append(f"Best average performer: {leader['company_name']} with {leader['stats']['average_views']:,} average views per sampled video.")
        cadence_leader = sorted(summaries, key=lambda s: s['stats']['video_cadence']['videos_per_month'], reverse=True)[0]
        insights.append(f"Most active publisher: {cadence_leader['company_name']} with {cadence_leader['stats']['video_cadence']['videos_per_month']} videos per month in the sample.")
        if leader['stats']['top_topics']:
            insights.append(f"Top content themes detected for {leader['company_name']}: {', '.join(leader['stats']['top_topics'][:3])}.")
    return insights


def build_recommendations(summaries):
    recommendations = []
    if not summaries:
        return recommendations
    leader = summaries[0]
    recommendations.append(f"Prioritize the strongest topic areas for {leader['company_name']} where engagement is highest.")
    for summary in summaries:
        if summary['stats']['content_gaps']:
            recommendations.append(f"Consider publishing new videos on {', '.join(summary['stats']['content_gaps'][:3])} to fill competitive gaps for {summary['company_name']}.")
    recommendations.append('Test weekly or bi-weekly publishing cadences if competitors are currently more active.')
    recommendations.append('Align titles and descriptions to top performing themes to improve discoverability and engagement.')
    return recommendations


def analyze_companies(companies):
    summaries = [compute_company_summary(name) for name in companies]
    global_topic_counts = Counter()
    for summary in summaries:
        global_topic_counts.update(summary['stats']['top_topics'])
    top_global_topics = [topic for topic, _ in global_topic_counts.most_common(10)]
    for summary in summaries:
        summary['stats']['content_gaps'] = [topic for topic in top_global_topics if topic not in summary['stats']['top_topics']][:6]
    ranked = score_companies(summaries)
    insights = build_insights(ranked)
    recommendations = build_recommendations(ranked)
    gap_insights = []
    for summary in ranked:
        if summary['stats']['content_gaps']:
            gap_insights.append(f"{summary['company_name']} lacks themes such as {', '.join(summary['stats']['content_gaps'][:3])}.")
    summary_text = {
        'company_count': len(summaries),
        'generated_at': datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC'),
    }
    return {
        'companies': summaries,
        'ranking': ranked,
        'insights': insights,
        'gap_insights': gap_insights,
        'recommendations': recommendations,
        'summary_text': summary_text,
    }


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return tf


def add_bullet_text(slide, left, top, width, height, lines, font_size=14):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = Pt(4)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            p.text = line
            first = False
        else:
            p = tf.add_paragraph()
            p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
    return tf


def add_table(slide, left, top, width, height, headers, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    return table


def set_slide_style(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND_COLOR


def add_slide_header(slide, title):
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT_COLOR
    header.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = HEADER_TEXT_COLOR
    p.alignment = PP_ALIGN.LEFT
    return title_box


def add_chart(slide, left, top, width, height, categories, series_data, title):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for legend, values in series_data.items():
        chart_data.add_series(legend, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2
    chart.legend.include_in_layout = False
    chart.chart_title.text_frame.text = title
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    return chart


def add_section_title(slide, text, top=Inches(1.2)):
    return add_text_box(slide, Inches(0.5), top, Inches(9), Inches(0.8), text, font_size=24, bold=True)


def create_pptx(analysis):
    presentation = Presentation()

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = BACKGROUND_COLOR
    cover.shapes.title.text = 'Competitive Video Intelligence Report'
    cover.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    cover.shapes.title.text_frame.paragraphs[0].font.bold = True
    cover.shapes.title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    subtitle = cover.shapes.placeholders[1]
    subtitle.text = f"Companies: {', '.join(c['company_name'] for c in analysis['companies'])}\nDate: {analysis['summary_text']['generated_at']}"
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)
    subtitle.text_frame.paragraphs[0].font.color.rgb = SECONDARY_TEXT_COLOR

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(summary_slide)
    add_slide_header(summary_slide, 'Executive Summary')
    add_bullet_text(summary_slide, Inches(0.5), Inches(1.4), Inches(9), Inches(4.8), analysis['insights'], font_size=16)

    overview_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(overview_slide)
    add_slide_header(overview_slide, 'Channel Overview')
    company_names = [summary['company_name'] for summary in analysis['companies']]
    add_chart(
        overview_slide,
        Inches(0.5), Inches(1.4), Inches(4.4), Inches(4.2),
        company_names,
        {
            'Avg Views': [summary['stats']['average_views'] for summary in analysis['companies']],
            'Subscribers': [summary['channel']['subscriber_count'] for summary in analysis['companies']],
        },
        'Average Views and Subscribers'
    )
    add_table(
        overview_slide,
        Inches(5.1), Inches(1.4), Inches(4.2), Inches(4.2),
        ['Company', 'Videos', 'Cadence', 'Videos / Month'],
        [[
            summary['company_name'],
            summary['stats']['total_videos'],
            summary['stats']['video_cadence']['cadence'],
            summary['stats']['video_cadence']['videos_per_month'],
        ] for summary in analysis['companies']]
    )

    engagement_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(engagement_slide)
    add_slide_header(engagement_slide, 'Engagement Performance')
    add_chart(
        engagement_slide,
        Inches(0.5), Inches(1.4), Inches(4.9), Inches(4.5),
        company_names,
        {
            'Avg Views': [summary['stats']['average_views'] for summary in analysis['companies']],
            'Avg Likes': [summary['stats']['average_likes'] for summary in analysis['companies']],
            'Avg Comments': [summary['stats']['average_comments'] for summary in analysis['companies']],
        },
        'Engagement by Company'
    )
    engagement_lines = []
    for summary in analysis['companies']:
        engagement_lines.append(
            f"{summary['company_name']}: top video '{summary['stats']['most_viewed_video']['title'] if summary['stats']['most_viewed_video'] else 'N/A'}'"
        )
        engagement_lines.append(
            f"Avg views {summary['stats']['average_views']:,}, likes {summary['stats']['average_likes']:,}, comments {summary['stats']['average_comments']:,}."
        )
        engagement_lines.append('')
    add_bullet_text(
        engagement_slide,
        Inches(5.5), Inches(1.4), Inches(3.8), Inches(4.5),
        engagement_lines,
        font_size=14
    )

    themes_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(themes_slide)
    add_slide_header(themes_slide, 'Content Themes & Gaps')
    add_bullet_text(
        themes_slide,
        Inches(0.5), Inches(1.4), Inches(4.7), Inches(4.8),
        [
            f"{summary['company_name']}: {', '.join(summary['stats']['top_topics'][:4])}."
            for summary in analysis['companies']
        ],
        font_size=14
    )
    add_bullet_text(
        themes_slide,
        Inches(5.3), Inches(1.4), Inches(4.0), Inches(4.8),
        [
            f"Gap: {summary['company_name']} lacks {', '.join(summary['stats']['content_gaps'][:3]) or 'key competitor themes'}."
            for summary in analysis['companies']
        ],
        font_size=14
    )

    recommendation_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(recommendation_slide)
    add_slide_header(recommendation_slide, 'Recommendations')
    add_bullet_text(recommendation_slide, Inches(0.5), Inches(1.4), Inches(9), Inches(5), analysis['recommendations'], font_size=16)

    ranking_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(ranking_slide)
    add_slide_header(ranking_slide, 'Ranking Summary')
    ranking_rows = [
        [summary['rank'], summary['company_name'], summary['score'], summary['stats']['average_views'], summary['stats']['average_likes']]
        for summary in analysis['ranking']
    ]
    add_table(
        ranking_slide,
        Inches(0.5), Inches(1.4), Inches(9), Inches(4.5),
        ['Rank', 'Company', 'Score', 'Avg Views', 'Avg Likes'],
        ranking_rows
    )

    appendix_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_slide_style(appendix_slide)
    add_slide_header(appendix_slide, 'Appendix: Top Sample Videos')
    appendix_lines = []
    for summary in analysis['companies']:
        appendix_lines.append(f"{summary['company_name']}:")
        for video in summary['stats']['top_videos'][:3]:
            appendix_lines.append(f" - {video['title']} ({video['view_count']:,} views)")
        appendix_lines.append('')
    add_bullet_text(appendix_slide, Inches(0.5), Inches(1.4), Inches(9), Inches(4.8), appendix_lines, font_size=12)

    return presentation


@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')


@app.route('/report', methods=['POST'])
def report():
    company_name = request.form.get('company_name', '').strip()
    competitors = [request.form.get(f'competitor_{i}', '').strip() for i in range(1, 5)]
    companies = [company_name] + [name for name in competitors if name]
    if not company_name:
        return render_template('index.html', error='Please enter your company name.', values=request.form)
    if len(companies) > 5:
        companies = companies[:5]
    analysis = analyze_companies(companies)
    return render_template('report.html', analysis=analysis, companies=companies)


@app.route('/download', methods=['POST'])
def download():
    companies = [request.form.get('company_name', '').strip()]
    companies += [request.form.get(f'competitor_{i}', '').strip() for i in range(1, 5)]
    companies = [name for name in companies if name]
    if not companies:
        return render_template('index.html', error='Please enter at least one company.')
    analysis = analyze_companies(companies)
    presentation = create_pptx(analysis)
    pptx_bytes = io.BytesIO()
    presentation.save(pptx_bytes)
    pptx_bytes.seek(0)
    filename = 'competitive_video_report.pptx'
    return send_file(
        pptx_bytes,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=False)


